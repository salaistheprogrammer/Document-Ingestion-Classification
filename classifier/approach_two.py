import os
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from pathlib import Path
from pdf2image import convert_from_path

def extract_pdf(file_path):
    lines = []

    images = convert_from_path(file_path, output_file=Path(__file__).parent)
    print(images)

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            words = []
            bboxes = []
            for word in page.extract_words():
                words.append(word['text'])
                bboxes.append([int(word['x0']), int(word['top']), int(word['x1']), int(word['bottom'])])
            lines.append({
                "image": images[i-1],
                "text": words,
                "bbox": bboxes,
                "page_num": i
            })
    return lines


def extract_docx(file_path):
    doc = Document(file_path)
    lines = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            lines.append({
                "text": text,
                "bbox": [0, i * 20, 200, (i + 1) * 20],  # Approximate vertical layout
                "page_num": 1
            })
    return lines


def extract_xlsx(file_path):
    wb = load_workbook(file_path, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                value = str(cell.value).strip() if cell.value else ""
                if value:
                    lines.append({
                        "text": value,
                        "bbox": [cell.column * 50, cell.row * 20, (cell.column + 1) * 50, (cell.row + 1) * 20],
                        "page_num": sheet.title
                    })
    return lines


def extract_pptx(file_path):
    prs = Presentation(file_path)
    lines = []
    for page_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append({
                        "text": text,
                        "page_num": page_num
                    })
    return lines


def extract_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".xlsx":
        return extract_xlsx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


# Example usage:
if __name__ == "__main__":
    # ip_path = r"C:\Hexaware\extractor\Standard deviation.docx"
    # extracted_docx = extract_from_file(ip_path)
    # print('-------------------------------docx---------------------------------')
    # for item in extracted_docx:
    #     print(item)

    # extracted_xlsx = extract_from_file(Path(__file__).parent / '1.xlsx')
    # print('-------------------------------xlsx---------------------------------')
    # for item in extracted_xlsx:
    #     print(item)

    # print('-------------------------------pptx---------------------------------')
    # extracted_pptx = extract_from_file(Path(__file__).parent / 'Team Outliers-SIH.pptx')
    # for item in extracted_pptx:
    #     print(item)

    print('--------------------------------pdf---------------------------------')
    extracted_pdf = extract_from_file(Path(__file__).parent / 'Document Ingestion & Classification.pdf')
    print(extracted_pdf)
